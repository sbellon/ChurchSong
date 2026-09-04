# SPDX-FileCopyrightText: 2024-2025 Stefan Bellon
#
# SPDX-License-Identifier: MIT

import logging
import os
import typing

import pptx
import pptx.shapes
import pptx.shapes.placeholder

from churchsong.powerpoint import PowerPointBase

if typing.TYPE_CHECKING:
    from churchsong.churchtools.events import Person
    from churchsong.configuration import Configuration


#######################################################################################
# Monkey patch python-pptx to accept MPO (multi picture object) JPEGs.
# Remove as soon as `ext_map` in `Image.ext()` contains "MPO" in pptx/parts/image.py
# https://github.com/scanny/python-pptx/issues/787

import pptx.parts.image
import pptx.util

_orig_ext = typing.cast('typing.Any', pptx.parts.image.Image.ext)._fget  # noqa: SLF001


@pptx.util.lazyproperty
def ext(self: pptx.parts.image.Image) -> str:
    return 'jpeg' if self._format == 'MPO' else _orig_ext(self)  # pyright: ignore # noqa: PGH003


pptx.parts.image.Image.ext = ext
#######################################################################################


logger = logging.getLogger(__name__)


class PowerPointServices(PowerPointBase):
    def __init__(self, config: Configuration) -> None:
        logger.info('Creating PowerPoint services slides')
        super().__init__(config, config.songbeamer.powerpoint.services.template_pptx)
        self._portraits_dir = config.songbeamer.powerpoint.services.portraits_dir

    def _insert_portrait(
        self, ph: pptx.shapes.placeholder.PicturePlaceholder, fullnames: str
    ) -> bool:
        portrait = self._portraits_dir / f'{fullnames}.jpeg'
        try:
            ph.insert_picture(  # pyright: ignore[reportUnknownMemberType]
                os.fspath(portrait)
            )
        except OSError as e:
            logger.error('Cannot embed portrait picture: %s', e)
            return False
        return True

    def create(self, service_leads: dict[str, set[Person]]) -> None:
        if not self._prs:
            return

        nobody = service_leads.get(str(None), set())
        slide_layout = self._prs.slide_layouts[0]
        slide = self._prs.slides.add_slide(slide_layout)
        for ph in slide.placeholders:
            base_placeholder = typing.cast(
                'pptx.shapes.placeholder.BasePlaceholder | None',
                getattr(ph, '_base_placeholder', None),
            )
            if not base_placeholder:
                logger.warning('Skipping unrecognized placeholder')
                continue
            service_name = base_placeholder.name
            sorted_persons = sorted(
                service_leads.get(service_name, nobody), key=lambda p: p.fullname
            )
            person_fullnames = ' + '.join(p.fullname for p in sorted_persons)
            person_shortnames = ' + '.join(p.shortname for p in sorted_persons)
            match ph:
                case pptx.shapes.placeholder.PicturePlaceholder():
                    logger.debug(
                        'Replacing image placeholder %s with %s',
                        service_name,
                        person_fullnames,
                    )
                    if not self._insert_portrait(ph, person_fullnames):
                        no_persons = ' + '.join(sorted(p.fullname for p in nobody))
                        if not self._insert_portrait(ph, no_persons):
                            logger.error(
                                'Leaving portrait placeholder %s empty', service_name
                            )
                case pptx.shapes.placeholder.SlidePlaceholder() if ph.has_text_frame:
                    logger.debug(
                        'Replacing text placeholder %s with %s',
                        service_name,
                        person_shortnames,
                    )
                    ph.text_frame.paragraphs[0].text = person_shortnames
                case _:
                    logger.warning(
                        'Skipping unsupported placeholder type %s',
                        ph.placeholder_format.type,
                    )
