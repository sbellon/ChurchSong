# SPDX-FileCopyrightText: 2024-2025 Stefan Bellon
#
# SPDX-License-Identifier: MIT

import datetime
import enum
import typing

import pptx.enum.dml
import pptx.shapes.graphfrm
import pptx.table
import pptx.text.text
import pptx.util

from churchsong.churchtools import RepeatId
from churchsong.configuration import CalendarSubtitleField
from churchsong.powerpoint import PowerPointBase

if typing.TYPE_CHECKING:
    from churchsong.churchtools import CalendarAppointmentBase
    from churchsong.configuration import Configuration


class TableType(enum.StrEnum):
    WEEKLY = 'Weekly Table'
    IRREGULAR = 'Irregular Table'


class TableFiller:
    def __init__(
        self,
        config: Configuration,
        table_type: TableType,
        one_week_later: datetime.datetime,
    ) -> None:
        self._log = config.log
        self._appointments_config = config.songbeamer.powerpoint.appointments
        self._table_type = table_type
        self._one_week_later = one_week_later
        self._subtitle_prio = (
            self._appointments_config.weekly.subtitle_priority
            if self._table_type == TableType.WEEKLY
            else self._appointments_config.irregular.subtitle_priority
        )
        self._table = None
        self._font = None
        self._total_rows = 0
        self._current_row = 0
        self._unset_table_warning = False

    @property
    def table_type(self) -> str:
        return self._table_type

    def set_table(self, table: pptx.table.Table) -> None:
        if self._table:
            self._log.warning('%s already set, not setting again', self._table_type)
        self._table = table
        self._total_rows = len(table.rows)
        self._current_row = 0
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        self._font = run.font
                        return

    def _set_font_properties(
        self,
        dst_font_element: pptx.text.text.Font,
        src_font_element: pptx.text.text.Font,
        scale: float = 1.0,
    ) -> None:
        for attr in ('name', 'bold', 'italic', 'underline', 'language_id'):
            if value := getattr(src_font_element, attr, None):
                setattr(dst_font_element, attr, value)
        if src_font_element.size:
            dst_font_element.size = pptx.util.Pt(src_font_element.size.pt * scale)

        match getattr(src_font_element.color, 'type', None):
            case pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                dst_font_element.color.rgb = getattr(
                    src_font_element.color, 'rgb', None
                )
            case pptx.enum.dml.MSO_COLOR_TYPE.SCHEME:
                dst_font_element.color.theme_color = getattr(
                    src_font_element.color, 'theme_color', None
                )
            case _:
                pass

    def _set_cell_text(
        self,
        cell: pptx.table._Cell,  # pyright: ignore[reportPrivateUsage]
        line1: str,
        line2: str | None = None,
    ) -> None:
        font_of_run = {
            idx: run.font for idx, run in enumerate(cell.text_frame.paragraphs[0].runs)
        }
        cell.text_frame.paragraphs[0].text = f'{line1}\v{line2}' if line2 else line1
        for idx, run in enumerate(cell.text_frame.paragraphs[0].runs):
            if font := font_of_run.get(idx) or self._font:
                self._set_font_properties(
                    run.font,
                    font,
                    scale=0.66 if idx > 0 and idx not in font_of_run else 1.0,
                )

    def _date_and_time(self, appt: CalendarAppointmentBase) -> str:
        # Use local timezone for display purposes.
        local_start = appt.start_date.astimezone()
        local_end = appt.end_date.astimezone()
        # Use 'weekly' format for appointments within the next week (on weekly
        # table as well as on irregular table), 'irregular' format for appointments
        # on irregular table, further away than a week.
        fmt_cfg = (
            self._appointments_config.weekly
            if appt.start_date < self._one_week_later
            else self._appointments_config.irregular
        )
        # Format date/time according to the following format strings with priority:
        #  1. multiday, if start != end
        #  2. allday, if marked allday
        #  3. regular (with time)
        if (local_start.month, local_start.day) != (local_end.month, local_end.day):
            return (
                f'{local_start:{fmt_cfg.multiday_datetime_format}}'
                ' - '
                f'{local_end:{fmt_cfg.multiday_datetime_format}}'
            )
        if appt.all_day:
            return f'{local_start:{fmt_cfg.allday_datetime_format}}'
        return f'{local_start:{fmt_cfg.regular_datetime_format}}'

    def _subtitle(self, appt: CalendarAppointmentBase) -> str:
        for subtitle in self._subtitle_prio:
            match subtitle:
                case CalendarSubtitleField.SUBTITLE if appt.subtitle:
                    return appt.subtitle
                case CalendarSubtitleField.DESCRIPTION if appt.description:
                    return appt.description
                case CalendarSubtitleField.LINK if appt.link:
                    return appt.link
                case CalendarSubtitleField.ADDRESS if appt.address:
                    city = f'{appt.address.zip or ""} {appt.address.city or ""}'.strip()
                    if address := (
                        ', '.join(
                            part
                            for part in [appt.address.name, appt.address.street, city]
                            if part
                        )
                    ):
                        return address
                case _:
                    pass
        return ''

    def add(self, appt: CalendarAppointmentBase) -> None:
        if not self._table:
            # Safeguard, no table registered.
            if not self._unset_table_warning:
                self._unset_table_warning = True
                self._log.warning(
                    '%s unset, ignoring all appointments', self._table_type
                )
            return
        if self._current_row >= self._total_rows:
            # All available table rows have been filled.
            self._log.info('%s is full, dropping "%s"', self._table_type, appt.title)
            return
        self._set_cell_text(
            self._table.cell(self._current_row, 0),
            self._date_and_time(appt),
        )
        self._set_cell_text(
            self._table.cell(self._current_row, 1),
            appt.title,
            self._subtitle(appt),
        )
        self._current_row += 1

    def fill(self) -> None:
        if not self._table:
            # Safeguard, no table registered.
            return
        for row in range(self._current_row, self._total_rows):
            self._set_cell_text(self._table.cell(row, 0), '')
            self._set_cell_text(self._table.cell(row, 1), '')


class PowerPointAppointments(PowerPointBase):
    def __init__(
        self, config: Configuration, event_start_date: datetime.datetime
    ) -> None:
        config.log.info('Creating PowerPoint appointments slides')
        super().__init__(
            config, config.songbeamer.powerpoint.appointments.template_pptx
        )
        self._one_week_later = event_start_date + datetime.timedelta(days=8)
        self._weekly_table = TableFiller(
            config=config,
            table_type=TableType.WEEKLY,
            one_week_later=self._one_week_later,
        )
        self._irregular_table = TableFiller(
            config=config,
            table_type=TableType.IRREGULAR,
            one_week_later=self._one_week_later,
        )

    def _setup_tables(self) -> None:
        if not self._prs:
            return

        # Walk through the slides and shapes and register the weekly table and the
        # irregular table for later filling with the appropriate appointments.
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if (
                    isinstance(shape, pptx.shapes.graphfrm.GraphicFrame)
                    and shape.has_table
                ):
                    match shape.name:
                        case self._weekly_table.table_type:
                            self._weekly_table.set_table(shape.table)
                        case self._irregular_table.table_type:
                            self._irregular_table.set_table(shape.table)
                        case _:
                            pass

    def create(self, appointments: typing.Iterable[CalendarAppointmentBase]) -> None:
        if not self._prs:
            return

        self._setup_tables()

        # Walk through the appointments and put them in the appropriate table.
        for appt in appointments:
            match (appt.repeat_id, appt.repeat_frequency):
                case (RepeatId.WEEKLY, 1) if appt.start_date < self._one_week_later:
                    self._weekly_table.add(appt)
                case (RepeatId.WEEKLY, 1):
                    pass  # ignore weekly appointments more than one week away
                case _:
                    self._irregular_table.add(appt)

        # Fill the tables to clean potential style templates in the cells.
        self._weekly_table.fill()
        self._irregular_table.fill()
