from __future__ import annotations

import os
import typing

import pptx
import pptx.exc
import pptx.shapes
import pptx.shapes.placeholder

if typing.TYPE_CHECKING:
    from churchsong.churchtools.events import Person
    from churchsong.configuration import Configuration


class PowerPoint:
    def __init__(self, config: Configuration) -> None:
        self._log = config.log
        self._portraits_dir = config.portraits_dir
        self._temp_dir = config.temp_dir
        self._template_pptx = config.template_pptx
        try:
            self._prs = pptx.Presentation(os.fspath(self._template_pptx))
        except pptx.exc.PackageNotFoundError as e:
            self._log.error(f'Cannot load PowerPoint template: {e}')
            self._prs = pptx.Presentation()

    def create(self, service_leads: dict[str, set[Person]]) -> None:
        self._log.info('Creating PowerPoint slide')
        if self._prs.core_properties.revision == 1:
            # Presentation is the fallback created one, just skip everything
            return
        slide_layout = self._prs.slide_layouts[0]
        slide = self._prs.slides.add_slide(slide_layout)
        for ph in slide.placeholders:
            service_name = ph._base_placeholder.name  # noqa: SLF001 # pyright: ignore[reportAttributeAccessIssue]
            sorted_persons = sorted(
                service_leads[service_name], key=lambda p: p.fullname
            )
            person_fullnames = ' + '.join(p.fullname for p in sorted_persons)
            person_shortnames = ' + '.join(p.shortname for p in sorted_persons)
            if isinstance(ph, pptx.shapes.placeholder.PicturePlaceholder):
                self._log.debug(
                    'Replacing image placeholder %s with %s',
                    service_name,
                    person_fullnames,
                )
                try:
                    ph.insert_picture(
                        os.fspath(self._portraits_dir / f'{person_fullnames}.jpeg')
                    )
                except FileNotFoundError as e:
                    self._log.error(f'Cannot embed portrait picture: {e}')
                    no_persons = ' + '.join(
                        sorted(p.fullname for p in service_leads[str(None)])
                    )
                    ph.insert_picture(
                        os.fspath(self._portraits_dir / f'{no_persons}.jpeg')
                    )
            elif (
                isinstance(ph, pptx.shapes.placeholder.SlidePlaceholder)
                and ph.has_text_frame
            ):
                self._log.debug(
                    'Replacing text placeholder %s with %s',
                    service_name,
                    person_shortnames,
                )
                ph.text_frame.paragraphs[0].text = person_shortnames
        _ = service_leads.pop(str(None), None)

    def save(self) -> None:
        self._prs.save(os.fspath(self._temp_dir / self._template_pptx.name))
