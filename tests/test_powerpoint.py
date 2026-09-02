# SPDX-FileCopyrightText: 2026 Stefan Bellon
#
# SPDX-License-Identifier: MIT

import base64
import copy
import datetime
import logging
import os
import typing

import pptx
import pptx.dml.color
import pptx.enum.dml
import pptx.enum.shapes
import pptx.oxml.ns
import pptx.shapes.graphfrm
import pptx.shapes.placeholder
import pptx.table
import pptx.util

from churchsong.churchtools import CalendarAppointmentBase
from churchsong.churchtools.events import Person
from churchsong.powerpoint.appointments import PowerPointAppointments
from churchsong.powerpoint.services import PowerPointServices
from tests.conftest import make_config

if typing.TYPE_CHECKING:
    import pathlib

    import pytest

    from churchsong.configuration import Configuration

# Smallest JPEG python-pptx accepts: a 1x1 pixel image.
JPEG_1PX = base64.b64decode(
    '/9j/4AAQSkZJRgABAQEAYABgAAD/2wBDAAgGBgcGBQgHBwcJCQgKDBQNDAsLDBkSEw8UHRof'
    'Hh0aHBwgJC4nICIsIxwcKDcpLDAxNDQ0Hyc5PTgyPC4zNDL/wAALCAABAAEBAREA/8QAFAAB'
    'AAAAAAAAAAAAAAAAAAAACf/EABQQAQAAAAAAAAAAAAAAAAAAAAD/2gAIAQEAAD8AKp//2Q=='
)


def make_services_template(
    path: pathlib.Path,
    *,
    text_services: tuple[str, str],
    picture_service: str,
    table_service: str | None = None,
) -> None:
    """Build a services template like a user would in PowerPoint.

    PowerPointServices instantiates slide layout 0 and matches placeholders to services
    by *name*, so rename layout 0's title/subtitle placeholders and insert a picture
    placeholder (borrowed from the stock 'Picture with Caption' layout) onto it.
    An optional table placeholder stands in for a placeholder type that cannot
    hold service staff at all.
    """
    prs = pptx.Presentation()
    layout = prs.slide_layouts[0]
    layout_placeholders = typing.cast(
        'typing.Iterable[pptx.shapes.placeholder.LayoutPlaceholder]',
        layout.placeholders,
    )
    title, subtitle = list(layout_placeholders)[:2]
    title.name = text_services[0]
    subtitle.name = text_services[1]
    picture_ph = next(
        ph
        for ph in prs.slide_layouts[8].placeholders
        if ph.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.PICTURE
    )
    element: typing.Any = copy.deepcopy(picture_ph._element)  # noqa: SLF001 # pyright: ignore[reportPrivateUsage]
    ph_element = element.find(f'.//{pptx.oxml.ns.qn("p:ph")}')
    assert ph_element is not None
    ph_element.set('idx', '13')  # unique placeholder idx within layout 0
    sp_tree = typing.cast('typing.Any', layout.shapes)._spTree  # noqa: SLF001
    sp_tree.append(element)
    picture_ph = next(
        ph for ph in layout.placeholders if ph.placeholder_format.idx == 13
    )
    picture_ph.name = picture_service
    if table_service:
        table_element: typing.Any = copy.deepcopy(element)
        table_ph_element = table_element.find(f'.//{pptx.oxml.ns.qn("p:ph")}')
        assert table_ph_element is not None
        table_ph_element.set('idx', '14')
        table_ph_element.set('type', 'tbl')
        sp_tree.append(table_element)
        table_ph = next(
            ph for ph in layout.placeholders if ph.placeholder_format.idx == 14
        )
        table_ph.name = table_service
    prs.save(os.fspath(path))


def make_appointments_template(
    path: pathlib.Path, *, weekly_rows: int, irregular_rows: int
) -> None:
    """Build an appointments template with the two named tables."""
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    def add_table(name: str, rows: int) -> None:
        frame = slide.shapes.add_table(
            rows,
            2,
            pptx.util.Cm(2),
            pptx.util.Cm(2),
            pptx.util.Cm(16),
            pptx.util.Cm(8),
        )
        frame.name = name
        for row in frame.table.rows:
            for cell in row.cells:
                cell.text = 'TEMPLATE'

    add_table('Weekly Table', weekly_rows)
    add_table('Irregular Table', irregular_rows)
    prs.save(os.fspath(path))


def read_tables(path: pathlib.Path) -> dict[str, list[list[str]]]:
    prs = pptx.Presentation(os.fspath(path))
    return {
        shape.name: [
            [cell.text_frame.text for cell in row.cells] for row in shape.table.rows
        ]
        for slide in prs.slides
        for shape in slide.shapes
        if isinstance(shape, pptx.shapes.graphfrm.GraphicFrame) and shape.has_table
    }


def make_appointment(  # noqa: PLR0913
    title: str,
    start: str,
    end: str | None = None,
    *,
    subtitle: str | None = None,
    description: str | None = None,
    link: str | None = None,
    address: dict[str, str | None] | None = None,
    all_day: bool = False,
    repeat_id: int = 0,
    repeat_frequency: int | None = None,
) -> CalendarAppointmentBase:
    return CalendarAppointmentBase.model_validate(
        {
            'title': title,
            'subtitle': subtitle,
            'description': description,
            'image': None,
            'link': link,
            'isInternal': False,
            'startDate': start,
            'endDate': end or start,
            'allDay': all_day,
            'repeatId': repeat_id,
            'repeatFrequency': repeat_frequency,
            'address': address,
        }
    )


def make_powerpoint_config(
    tmp_path: pathlib.Path, section: str, template: pathlib.Path
) -> Configuration:
    powerpoint: dict[str, dict[str, object]] = {
        section: {'template_pptx': str(template)}
    }
    if section == 'Services':
        powerpoint[section]['portraits_dir'] = str(tmp_path / 'portraits')
    if section == 'Appointments':
        powerpoint[section].update(
            {
                'Weekly': {
                    'regular_datetime_format': 'W %d.%m. %H:%M',
                    'allday_datetime_format': 'W %d.%m.',
                    'multiday_datetime_format': 'W %d.%m.',
                },
                'Irregular': {
                    'regular_datetime_format': 'I %d.%m. %H:%M',
                    'allday_datetime_format': 'I %d.%m.',
                    'multiday_datetime_format': 'I %d.%m.',
                },
            }
        )
    return make_config(
        output_dir=str(tmp_path / 'output'),
        songbeamer={'PowerPoint': powerpoint},
    )


def test_services_fills_text_and_picture_placeholders(
    tmp_path: pathlib.Path,
) -> None:
    template = tmp_path / 'services.pptx'
    make_services_template(
        template, text_services=('Preaching', 'Music'), picture_service='Welcome'
    )
    portraits = tmp_path / 'portraits'
    portraits.mkdir()
    (portraits / 'Jane Doe.jpeg').write_bytes(JPEG_1PX)
    (tmp_path / 'output').mkdir()
    config = make_powerpoint_config(tmp_path, 'Services', template)

    service_leads = {
        'Preaching': {Person('Jane Doe', 'Jane')},
        'Music': {Person('Bo Li', 'Bo'), Person('Alex Roe', 'Alex')},
        'Welcome': {Person('Jane Doe', 'Jane')},
        str(None): {Person('Nobody', 'Nobody')},
    }
    powerpoint = PowerPointServices(config)
    powerpoint.create(service_leads)
    powerpoint.save()

    assert str(None) not in service_leads  # popped after use

    result = pptx.Presentation(os.fspath(tmp_path / 'output' / 'services.pptx'))
    (slide,) = result.slides
    # Slide placeholders inherit by idx from the layout placeholders whose
    # names select the service: idx 0 was named 'Preaching', 1 'Music'.
    texts = {
        shape.placeholder_format.idx: typing.cast(
            'pptx.shapes.placeholder.SlidePlaceholder', shape
        ).text_frame.text
        for shape in slide.placeholders
        if not isinstance(shape, pptx.shapes.placeholder.PlaceholderPicture)
    }
    assert texts == {
        0: 'Jane',
        1: 'Alex + Bo',  # sorted by fullname, joined shortnames
    }
    (picture,) = (
        shape
        for shape in slide.placeholders
        if isinstance(shape, pptx.shapes.placeholder.PlaceholderPicture)
    )
    assert picture.image.blob == JPEG_1PX


def test_services_falls_back_to_nobody_portrait(tmp_path: pathlib.Path) -> None:
    template = tmp_path / 'services.pptx'
    make_services_template(
        template, text_services=('Preaching', 'Music'), picture_service='Welcome'
    )
    portraits = tmp_path / 'portraits'
    portraits.mkdir()
    (portraits / 'Nobody.jpeg').write_bytes(JPEG_1PX)  # no 'Jane Doe.jpeg'
    (tmp_path / 'output').mkdir()
    config = make_powerpoint_config(tmp_path, 'Services', template)

    jane = Person('Jane Doe', 'Jane')
    powerpoint = PowerPointServices(config)
    powerpoint.create(
        {
            'Preaching': {jane},
            'Music': {jane},
            'Welcome': {jane},
            str(None): {Person('Nobody', 'Nobody')},
        }
    )
    powerpoint.save()

    result = pptx.Presentation(os.fspath(tmp_path / 'output' / 'services.pptx'))
    (picture,) = (
        shape
        for shape in result.slides[0].placeholders
        if isinstance(shape, pptx.shapes.placeholder.PlaceholderPicture)
    )
    assert picture.image.blob == JPEG_1PX


def test_appointments_distributes_and_formats_appointments(
    tmp_path: pathlib.Path,
) -> None:
    template = tmp_path / 'appointments.pptx'
    make_appointments_template(template, weekly_rows=2, irregular_rows=5)
    (tmp_path / 'output').mkdir()
    config = make_powerpoint_config(tmp_path, 'Appointments', template)
    event_date = datetime.datetime(2026, 8, 23, 10, 0, tzinfo=datetime.UTC)

    def local(day: int, hour: int, month: int = 8) -> datetime.datetime:
        return datetime.datetime(
            2026, month, day, hour, tzinfo=datetime.UTC
        ).astimezone()

    appointments = [
        # Weekly repetitions within the next week go to the weekly table ...
        make_appointment(
            'Prayer Meeting',
            '2026-08-25T09:00:00Z',
            subtitle='Chapel',
            repeat_id=7,
            repeat_frequency=1,
        ),
        make_appointment(
            'Youth Group', '2026-08-26T18:00:00Z', repeat_id=7, repeat_frequency=1
        ),
        # ... until the table is full ...
        make_appointment(
            'Overflow', '2026-08-27T18:00:00Z', repeat_id=7, repeat_frequency=1
        ),
        # ... and are ignored entirely when more than a week away.
        make_appointment(
            'Far Weekly', '2026-09-15T18:00:00Z', repeat_id=7, repeat_frequency=1
        ),
        # Everything else goes to the irregular table.
        make_appointment('Festival', '2026-08-24T10:00:00Z', all_day=True),
        make_appointment(
            'Concert', '2026-09-05T14:00:00Z', description='With the choir'
        ),
        make_appointment(
            'Retreat',
            '2026-09-10T10:00:00Z',
            '2026-09-12T10:00:00Z',
            link='https://retreat.test',
        ),
        make_appointment(
            'Bake Sale',
            '2026-09-19T09:00:00Z',
            address={
                'name': 'Church Hall',
                'street': 'Main St 1',
                'zip': '12345',
                'city': 'Springfield',
            },
        ),
    ]
    powerpoint = PowerPointAppointments(config, event_date)
    powerpoint.create(appointments)
    powerpoint.save()

    tables = read_tables(tmp_path / 'output' / 'appointments.pptx')
    assert tables['Weekly Table'] == [
        [f'{local(25, 9):W %d.%m. %H:%M}', 'Prayer Meeting\vChapel'],
        [f'{local(26, 18):W %d.%m. %H:%M}', 'Youth Group'],
    ]
    assert tables['Irregular Table'] == [
        # All-day within the next week: weekly all-day format.
        [f'{local(24, 10):W %d.%m.}', 'Festival'],
        # Subtitle falls back to description, link, then address.
        [f'{local(5, 14, month=9):I %d.%m. %H:%M}', 'Concert\vWith the choir'],
        [
            f'{local(10, 10, month=9):I %d.%m.} - {local(12, 10, month=9):I %d.%m.}',
            'Retreat\vhttps://retreat.test',
        ],
        [
            f'{local(19, 9, month=9):I %d.%m. %H:%M}',
            'Bake Sale\vChurch Hall, Main St 1, 12345 Springfield',
        ],
        # Unused rows are cleared of template content.
        ['', ''],
    ]


def test_appointments_without_matching_tables_are_dropped(
    tmp_path: pathlib.Path,
) -> None:
    template = tmp_path / 'appointments.pptx'
    # Template with a weekly table only: irregular appointments have no home.
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_table(
        1,
        2,
        pptx.util.Cm(2),
        pptx.util.Cm(2),
        pptx.util.Cm(16),
        pptx.util.Cm(2),
    )
    frame.name = 'Weekly Table'
    prs.save(os.fspath(template))
    (tmp_path / 'output').mkdir()
    config = make_powerpoint_config(tmp_path, 'Appointments', template)
    event_date = datetime.datetime(2026, 8, 23, 10, 0, tzinfo=datetime.UTC)

    powerpoint = PowerPointAppointments(config, event_date)
    powerpoint.create([make_appointment('Concert', '2026-09-05T14:00:00Z')])
    powerpoint.save()

    tables = read_tables(tmp_path / 'output' / 'appointments.pptx')
    assert 'Irregular Table' not in tables


def test_missing_template_configuration_skips_powerpoint(
    tmp_path: pathlib.Path,
) -> None:
    config = make_config(output_dir=str(tmp_path))
    powerpoint = PowerPointServices(config)
    powerpoint.create({})
    powerpoint.save()
    assert list(tmp_path.iterdir()) == []


def test_corrupt_template_skips_powerpoint(tmp_path: pathlib.Path) -> None:
    template = tmp_path / 'services.pptx'
    template.write_bytes(b'this is not a pptx file')
    (tmp_path / 'output').mkdir()
    config = make_powerpoint_config(tmp_path, 'Services', template)
    powerpoint = PowerPointServices(config)
    powerpoint.create({})
    powerpoint.save()
    assert list((tmp_path / 'output').iterdir()) == []


def make_odd_appointments_template(path: pathlib.Path) -> None:
    """A template with a duplicate weekly table and an unrelated table."""
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for name in ('Weekly Table', 'Weekly Table', 'Other Table'):
        frame = slide.shapes.add_table(
            1, 2, pptx.util.Cm(2), pptx.util.Cm(2), pptx.util.Cm(16), pptx.util.Cm(2)
        )
        frame.name = name
        for cell in frame.table.rows[0].cells:
            cell.text = 'TEMPLATE'
    prs.save(os.fspath(path))


def make_styled_appointments_template(
    path: pathlib.Path, *, theme_color: bool = False
) -> None:
    """A template whose table text carries explicit character formatting."""
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_table(
        1, 2, pptx.util.Cm(2), pptx.util.Cm(2), pptx.util.Cm(16), pptx.util.Cm(2)
    )
    frame.name = 'Weekly Table'
    for cell in frame.table.rows[0].cells:
        cell.text = 'TEMPLATE'
        font = cell.text_frame.paragraphs[0].runs[0].font
        font.bold = True
        font.size = pptx.util.Pt(18)
        if theme_color:
            font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.ACCENT_1
        else:
            font.color.rgb = pptx.dml.color.RGBColor(0xFF, 0x00, 0x00)
    prs.save(os.fspath(path))


def read_first_table(path: pathlib.Path) -> pptx.table.Table:
    prs = pptx.Presentation(os.fspath(path))
    return next(
        shape.table
        for slide in prs.slides
        for shape in slide.shapes
        if isinstance(shape, pptx.shapes.graphfrm.GraphicFrame) and shape.has_table
    )


def create_appointments(tmp_path: pathlib.Path, template: pathlib.Path) -> pathlib.Path:
    """Render one weekly appointment into the template and save the result."""
    (tmp_path / 'output').mkdir(exist_ok=True)
    config = make_powerpoint_config(tmp_path, 'Appointments', template)
    powerpoint = PowerPointAppointments(
        config, datetime.datetime(2026, 8, 23, 10, 0, tzinfo=datetime.UTC)
    )
    powerpoint.create(
        [
            make_appointment(
                'Prayer Meeting',
                '2026-08-25T09:00:00Z',
                subtitle='Chapel',
                repeat_id=7,
                repeat_frequency=1,
            )
        ]
    )
    powerpoint.save()
    return tmp_path / 'output' / 'appointments.pptx'


def test_missing_appointments_template_skips_powerpoint(
    tmp_path: pathlib.Path,
) -> None:
    config = make_config(output_dir=str(tmp_path))
    powerpoint = PowerPointAppointments(
        config, datetime.datetime(2026, 8, 23, 10, 0, tzinfo=datetime.UTC)
    )
    powerpoint.create([make_appointment('Concert', '2026-09-05T14:00:00Z')])
    powerpoint.save()
    assert list(tmp_path.iterdir()) == []


def test_duplicate_table_is_ignored_and_unknown_table_untouched(
    tmp_path: pathlib.Path, caplog: pytest.LogCaptureFixture
) -> None:
    template = tmp_path / 'appointments.pptx'
    make_odd_appointments_template(template)
    with caplog.at_level(logging.WARNING):
        output = create_appointments(tmp_path, template)
    assert 'Weekly Table already set' in caplog.text
    prs = pptx.Presentation(os.fspath(output))
    weekly = [
        [[cell.text_frame.text for cell in row.cells] for row in shape.table.rows]
        for slide in prs.slides
        for shape in slide.shapes
        if isinstance(shape, pptx.shapes.graphfrm.GraphicFrame)
        and shape.has_table
        and shape.name == 'Weekly Table'
    ]
    # The first of the two weekly tables wins, the duplicate stays untouched ...
    assert weekly[0][0][1] == 'Prayer Meeting\vChapel'
    assert weekly[1][0] == ['TEMPLATE', 'TEMPLATE']
    # ... and so does the table nobody asked for.
    assert read_tables(output)['Other Table'] == [['TEMPLATE', 'TEMPLATE']]


def test_appointment_text_inherits_the_template_font(tmp_path: pathlib.Path) -> None:
    template = tmp_path / 'appointments.pptx'
    make_styled_appointments_template(template)
    table = read_first_table(create_appointments(tmp_path, template))
    title_run, subtitle_run = table.cell(0, 1).text_frame.paragraphs[0].runs
    assert title_run.font.bold
    assert title_run.font.size == pptx.util.Pt(18)
    assert (
        title_run.font.color.rgb  # pyright: ignore[reportUnknownMemberType]
        == pptx.dml.color.RGBColor(0xFF, 0x00, 0x00)
    )
    # The subtitle is added by us and rendered smaller than the title.
    assert subtitle_run.font.size is not None
    assert title_run.font.size is not None
    assert subtitle_run.font.size < title_run.font.size


def test_appointment_text_inherits_a_theme_colored_template_font(
    tmp_path: pathlib.Path,
) -> None:
    template = tmp_path / 'appointments.pptx'
    make_styled_appointments_template(template, theme_color=True)
    table = read_first_table(create_appointments(tmp_path, template))
    title_run = table.cell(0, 1).text_frame.paragraphs[0].runs[0]
    assert (
        title_run.font.color.theme_color  # pyright: ignore[reportUnknownMemberType]
        == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_1
    )


def test_services_warns_about_an_unsupported_placeholder(
    tmp_path: pathlib.Path, caplog: pytest.LogCaptureFixture
) -> None:
    template = tmp_path / 'services.pptx'
    make_services_template(
        template,
        text_services=('Preaching', 'Music'),
        picture_service='Welcome',
        table_service='Coffee',  # a table cannot hold service staff
    )
    portraits = tmp_path / 'portraits'
    portraits.mkdir()
    (portraits / 'Jane Doe.jpeg').write_bytes(JPEG_1PX)
    (tmp_path / 'output').mkdir()
    config = make_powerpoint_config(tmp_path, 'Services', template)

    jane = {Person('Jane Doe', 'Jane')}
    powerpoint = PowerPointServices(config)
    with caplog.at_level(logging.WARNING):
        powerpoint.create(
            {
                'Preaching': jane,
                'Music': jane,
                'Welcome': jane,
                'Coffee': jane,
                str(None): {Person('Nobody', 'Nobody')},
            }
        )
    powerpoint.save()
    assert 'Skipping unsupported placeholder type' in caplog.text
    # The rest of the slide is still filled in.
    result = pptx.Presentation(os.fspath(tmp_path / 'output' / 'services.pptx'))
    texts = [
        shape.text_frame.text
        for shape in result.slides[0].placeholders
        if isinstance(shape, pptx.shapes.placeholder.SlidePlaceholder)
    ]
    assert 'Jane' in texts
