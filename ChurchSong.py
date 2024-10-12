#!/usr/bin/env python

from __future__ import annotations

import argparse
import configparser
import datetime
import io
import logging
import logging.handlers
import os
import pathlib
import subprocess
import sys
import typing
import zipfile
from collections import defaultdict

import marshmallow
import marshmallow_dataclass
import pptx
import pptx.shapes
import pptx.shapes.placeholder
import requests


class Configuration:
    def __init__(self, ini_file: pathlib.Path) -> None:
        self._log = logging.getLogger(__name__)
        self._log.setLevel(logging.INFO)
        log_formatter = logging.Formatter(
            '%(asctime)s - %(levelname)-8s - %(message)s', datefmt='%Y-%m-%d %H:%M:%S'
        )

        # Log to stderr before we have the log_file name from the .ini file.
        log_to_stderr = logging.StreamHandler(sys.stderr)
        log_to_stderr.setFormatter(log_formatter)
        self._log.addHandler(log_to_stderr)

        # Read the configuration .ini file.
        self._config = configparser.RawConfigParser(
            interpolation=configparser.ExtendedInterpolation(),
        )
        self._config.optionxform = lambda optionstr: optionstr
        self._config.read(ini_file, encoding='utf-8')

        # Switch to configured logging.
        self._log.setLevel(self.log_level)
        log_to_file = logging.handlers.RotatingFileHandler(
            self.log_file, maxBytes=5 * 1024 * 1024, backupCount=7
        )
        log_to_file.setFormatter(log_formatter)
        self._log.addHandler(log_to_file)
        self._log.removeHandler(log_to_stderr)

    def _expand_vars(self, section: str, option: str) -> str:
        return self._config.get(section, option, vars=dict(os.environ))

    @property
    def log(self) -> logging.Logger:
        return self._log

    @property
    def log_level(self) -> str:
        return self._expand_vars('General', 'log_level')

    @property
    def log_file(self) -> pathlib.Path:
        filename = pathlib.Path(self._expand_vars('General', 'log_file'))
        filename.parent.mkdir(parents=True, exist_ok=True)
        return filename

    @property
    def base_url(self) -> str:
        return self._expand_vars('ChurchTools.Settings', 'base_url')

    @property
    def login_token(self) -> str:
        return self._expand_vars('ChurchTools.Settings', 'login_token')

    @property
    def person_dict(self) -> dict[str, str]:
        return dict(self._config.items('ChurchTools.Replacements'))

    @property
    def template_pptx(self) -> pathlib.Path:
        return pathlib.Path(self._expand_vars('SongBeamer.Settings', 'template_pptx'))

    @property
    def portraits_dir(self) -> pathlib.Path:
        return pathlib.Path(self._expand_vars('SongBeamer.Settings', 'portraits_dir'))

    @property
    def temp_dir(self) -> pathlib.Path:
        directory = pathlib.Path(self._expand_vars('SongBeamer.Settings', 'temp_dir'))
        directory.mkdir(parents=True, exist_ok=True)
        return directory

    @property
    def replacements(self) -> list[tuple[str, str]]:
        return [
            (key, val.replace('\\n', '\n').replace('\\r', '\r').replace('\\', '/'))
            for key, val in self._config.items(
                'SongBeamer.Replacements',
                vars=dict(os.environ),
            )
        ]


@typing.dataclass_transform()
def deserialize(cls: type) -> type:
    class Meta:
        unknown = marshmallow.EXCLUDE

    cls.Meta = Meta
    return marshmallow_dataclass.dataclass(cls)


SchemaType = typing.ClassVar[type[marshmallow.Schema]]  # for pylance


@deserialize
class Service:
    Schema: SchemaType  # for pylance
    id: int
    name: str | None


@deserialize
class ServicesData:
    Schema: SchemaType  # for pylance
    data: list[Service]


@deserialize
class EventShort:
    Schema: SchemaType  # for pylance
    id: int
    startDate: datetime.datetime  # noqa: N815


@deserialize
class EventsData:
    Schema: SchemaType  # for pylance
    data: list[EventShort]


@deserialize
class EventService:
    Schema: SchemaType  # for pylance
    serviceId: int  # noqa: N815
    name: str | None


@deserialize
class EventFull:
    Schema: SchemaType  # for pylance
    id: int
    eventServices: list[EventService]  # noqa: N815


@deserialize
class EventFullData:
    Schema: SchemaType  # for pylance
    data: EventFull


@deserialize
class EventAgenda:
    Schema: SchemaType  # for pylance
    id: int


@deserialize
class EventAgendaData:
    Schema: SchemaType  # for pylance
    data: EventAgenda


@deserialize
class AgendaExport:
    Schema: SchemaType  # for pylance
    url: str


@deserialize
class AgendaExportData:
    Schema: SchemaType  # for pylance
    data: AgendaExport


class ChurchTools:
    def __init__(self, config: Configuration) -> None:
        self._log = config.log
        self._base_url = config.base_url
        self._login_token = config.login_token
        self._person_dict = config.person_dict
        self._temp_dir = config.temp_dir

    def _headers(self) -> dict[str, str]:
        return {
            'Accept': 'application/json',
            'Authorization': f'Login {self._login_token}',
        }

    def _request(
        self,
        method: str,
        url: str,
        params: dict[str, str] | None = None,
    ) -> requests.Response:
        self._log.debug(
            'Request %s %s%s with params=%s', method, self._base_url, url, params
        )
        r = requests.request(
            method,
            f'{self._base_url}{url}',
            headers=self._headers(),
            params=params,
            timeout=None,  # noqa: S113
        )
        self._log.debug('Response is %s %s', r.status_code, r.reason)
        r.raise_for_status()
        return r

    def _get(
        self,
        url: str,
        params: dict[str, str] | None = None,
    ) -> requests.Response:
        return self._request('GET', url, params)

    def _post(
        self,
        url: str,
        params: dict[str, str] | None = None,
    ) -> requests.Response:
        return self._request('POST', url, params)

    def _get_services(self) -> list[Service]:
        r = self._get('/api/services')
        result = ServicesData.Schema().load(r.json())
        assert isinstance(result, ServicesData)
        return result.data

    def _get_events(self, from_date: datetime.date | None = None) -> list[EventShort]:
        r = self._get(
            '/api/events', params={'from': f'{from_date}'} if from_date else None
        )
        result = EventsData.Schema().load(r.json())
        assert isinstance(result, EventsData)
        return result.data

    def _get_next_event(self, from_date: datetime.date | None = None) -> EventShort:
        try:
            return self._get_events(from_date)[0]
        except IndexError:
            err_msg = 'No events present{} in ChurchTools'.format(
                f' after {from_date}' if from_date else ''
            )
            self._log.error(err_msg)
            sys.stderr.write(f'{err_msg}\n')
            sys.exit(1)

    def _get_event(self, event_id: int) -> EventFull:
        r = self._get(f'/api/events/{event_id}')
        result = EventFullData.Schema().load(r.json())
        assert isinstance(result, EventFullData)
        return result.data

    def _get_event_agenda(self, event_id: int) -> EventAgenda:
        r = self._get(f'/api/events/{event_id}/agenda')
        result = EventAgendaData.Schema().load(r.json())
        assert isinstance(result, EventAgendaData)
        return result.data

    def _get_agenda_export(self, agenda_id: int) -> AgendaExport:
        r = self._post(
            f'/api/agendas/{agenda_id}/export',
            params={
                'target': 'SONG_BEAMER',
                'exportSongs': 'true',
                'appendArrangement': 'false',
                'withCategory': 'false',
            },
        )
        result = AgendaExportData.Schema().load(r.json())
        assert isinstance(result, AgendaExportData)
        return result.data

    def get_service_leads(
        self, from_date: datetime.date | None = None
    ) -> defaultdict[str, str]:
        self._log.info('Fetching service teams')
        services = self._get_services()
        next_event = self._get_next_event(from_date)
        event = self._get_event(next_event.id)
        # Initialize the "None" person for all services.
        service_leads = defaultdict(lambda: self._person_dict.get(str(None), str(None)))
        # Update with the actual persons of the eventservice.
        service_leads.update(
            {
                service.name: self._person_dict.get(
                    str(eventservice.name),
                    str(eventservice.name),
                )
                for eventservice in event.eventServices
                for service in services
                if eventservice.serviceId == service.id
            },
        )
        return service_leads

    def get_url_for_songbeamer_agenda(
        self, from_date: datetime.date | None = None
    ) -> str:
        self._log.info('Fetching SongBeamer export URL')
        next_event = self._get_next_event(from_date)
        try:
            agenda = self._get_event_agenda(next_event.id)
        except requests.HTTPError as e:
            if e.response.status_code == requests.codes['not_found']:
                date = next_event.startDate.date()
                err_msg = f'No event agenda present for {date} in ChurchTools'
                self._log.error(err_msg)
                sys.stderr.write(f'{err_msg}\n')
                sys.exit(1)
            raise
        return self._get_agenda_export(agenda.id).url

    def download_and_extract_agenda_zip(self, url: str) -> None:
        self._log.info('Downloading and extracting SongBeamer export')
        r = self._get(url)
        assert isinstance(r.content, bytes)
        buf = io.BytesIO(r.content)
        zipfile.ZipFile(buf, mode='r').extractall(path=self._temp_dir)


class PowerPoint:
    def __init__(self, config: Configuration) -> None:
        self._log = config.log
        self._portraits_dir = config.portraits_dir
        self._temp_dir = config.temp_dir
        self._template_pptx = config.template_pptx
        self._prs = pptx.Presentation(os.fspath(self._template_pptx))

    def create(self, service_leads: dict[str, str]) -> None:
        self._log.info('Creating PowerPoint slide')
        slide_layout = self._prs.slide_layouts[0]
        slide = self._prs.slides.add_slide(slide_layout)
        for ph in slide.placeholders:
            service_name = ph._base_placeholder.name  # noqa: SLF001 # pyright: ignore[reportAttributeAccessIssue]
            person_name = service_leads[service_name]
            if isinstance(ph, pptx.shapes.placeholder.PicturePlaceholder):
                self._log.debug(
                    'Replacing image placeholder %s with %s', service_name, person_name
                )
                ph.insert_picture(
                    os.fspath(self._portraits_dir / f'{person_name}.jpeg')
                )
            elif (
                isinstance(ph, pptx.shapes.placeholder.SlidePlaceholder)
                and ph.has_text_frame
            ):
                self._log.debug(
                    'Replacing text placeholder %s with %s', service_name, person_name
                )
                ph.text_frame.paragraphs[0].text = person_name.split(' ')[0]

    def save(self) -> None:
        self._prs.save(os.fspath(self._temp_dir / self._template_pptx.name))


class SongBeamer:
    def __init__(self, config: Configuration) -> None:
        self._log = config.log
        self._temp_dir = config.temp_dir.resolve()
        self._schedule_filepath = self._temp_dir / 'Schedule.col'
        self._replacements = config.replacements

    def modify_and_save_agenda(self) -> None:
        self._log.info('Modifying SongBeamer schedule')
        with self._schedule_filepath.open(mode='r', encoding='utf-8') as fd:
            content = fd.read()
        for search, replace in self._replacements:
            content = content.replace(search, replace)
        with self._schedule_filepath.open(mode='w', encoding='utf-8') as fd:
            fd.write(content)

    def launch(self) -> None:
        self._log.info('Launching SongBeamer instance')
        subprocess.run(
            [os.environ.get('COMSPEC', 'cmd'), '/C', 'start Schedule.col'],
            check=True,
            cwd=self._temp_dir,
        )


def main() -> None:
    config = Configuration(pathlib.Path(__file__).with_suffix('.ini'))

    config.log.debug('Parsing command line with args: %s', sys.argv)
    parser = argparse.ArgumentParser(
        prog='ChurchSong',
        description='Download ChurchTools event agenda and import into SongBeamer.',
    )
    parser.add_argument(
        'from_date',
        metavar='FROM_DATE',
        type=datetime.date.fromisoformat,
        nargs='?',
        help='search in ChurchTools for next event starting at FROM_DATE (YYYY-MM-DD)',
    )
    args = parser.parse_args()

    config.log.info('Starting ChurchSong with FROM_DATE=%s', args.from_date)
    try:
        ct = ChurchTools(config)
        service_leads = ct.get_service_leads(args.from_date)

        pp = PowerPoint(config)
        pp.create(service_leads)
        pp.save()

        ct.download_and_extract_agenda_zip(
            ct.get_url_for_songbeamer_agenda(args.from_date),
        )

        sb = SongBeamer(config)
        sb.modify_and_save_agenda()
        sb.launch()
    except Exception as e:
        config.log.fatal(e, exc_info=True)
        raise


if __name__ == '__main__':
    main()
