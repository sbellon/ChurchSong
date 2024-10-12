from __future__ import annotations

import os
import typing

import pptx
import pptx.shapes
import pptx.shapes.placeholder

if typing.TYPE_CHECKING:
    from configuration import Configuration


class PowerPoint:
    def __init__(self, config: Configuration) -> None:
        self._log = config.log
        self._portraits_dir = config.portraits_dir
        self._temp_dir = config.temp_dir
        self._template_pptx = config.template_pptx
        self._prs = pptx.Presentation(os.fspath(self._template_pptx))

    def create(self, service_leads: dict[str, str]) -> None:
        self._log.info('Creating PowerPoint slide')
        slide_layout = self._prs.slide_layouts[0]
        slide = self._prs.slides.add_slide(slide_layout)
        for ph in slide.placeholders:
            service_name = ph._base_placeholder.name  # noqa: SLF001 # pyright: ignore[reportAttributeAccessIssue]
            person_name = service_leads[service_name]
            if isinstance(ph, pptx.shapes.placeholder.PicturePlaceholder):
                self._log.debug(
                    'Replacing image placeholder %s with %s', service_name, person_name
                )
                ph.insert_picture(
                    os.fspath(self._portraits_dir / f'{person_name}.jpeg')
                )
            elif (
                isinstance(ph, pptx.shapes.placeholder.SlidePlaceholder)
                and ph.has_text_frame
            ):
                self._log.debug(
                    'Replacing text placeholder %s with %s', service_name, person_name
                )
                ph.text_frame.paragraphs[0].text = person_name.split(' ')[0]

    def save(self) -> None:
        self._prs.save(os.fspath(self._temp_dir / self._template_pptx.name))
